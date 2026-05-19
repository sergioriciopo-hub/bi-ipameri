# -*- coding: utf-8 -*-
"""
pptx_report.py — Gerador de Apresentacao PPTX para BI Ipameri/Buriti
Chamado pelo Streamlit via: gerar_pptx(D, d0, d1, CFG) -> bytes

Slide size: 20" x 11.25" (widescreen)
Graficos: matplotlib (sem kaleido)
"""

import io
import warnings
warnings.filterwarnings("ignore")

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Dimensoes do slide ────────────────────────────────────────────────────────
W = Inches(20)
H = Inches(11.25)

MESES_PT = {1:"Jan",2:"Fev",3:"Mar",4:"Abr",5:"Mai",6:"Jun",
            7:"Jul",8:"Ago",9:"Set",10:"Out",11:"Nov",12:"Dez"}

# ── Configuracoes por cliente ─────────────────────────────────────────────────
CFG_IPAMERI = dict(
    nome="Águas de Ipameri",
    sigla="IPAMERI",
    bi_label="Business Intelligence",
    cor_principal=RGBColor(0x1A, 0x6F, 0xAD),
    cor_escura=RGBColor(0x0D, 0x3B, 0x5E),
    cor_clara=RGBColor(0x5B, 0x8F, 0xB8),
    cor_hex="#1A6FAD",
    cor_esc_hex="#0D3B5E",
    cor_clr_hex="#5B8FB8",
    url_outro="https://bi-buriti-alegre.streamlit.app",
    label_outro="BI Buriti Alegre",
)

CFG_BURITI = dict(
    nome="Buriti Alegre Ambiental",
    sigla="BURITI",
    bi_label="Business Intelligence",
    cor_principal=RGBColor(0x2E, 0x7D, 0x32),
    cor_escura=RGBColor(0x1B, 0x5E, 0x20),
    cor_clara=RGBColor(0x66, 0xBB, 0x6A),
    cor_hex="#2E7D32",
    cor_esc_hex="#1B5E20",
    cor_clr_hex="#66BB6A",
    url_outro="https://bi-ipameri.streamlit.app",
    label_outro="BI Ipameri",
)

# Cores fixas
BRANCO   = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_BG = RGBColor(0xF4, 0xF6, 0xF8)
CINZA_TX = RGBColor(0x55, 0x65, 0x75)
VERMELHO = RGBColor(0xC0, 0x39, 0x2B)
VERDE_C  = RGBColor(0x27, 0xAE, 0x60)
LARANJA  = RGBColor(0xE6, 0x7E, 0x22)


# ══════════════════════════════════════════════════════════════════════════════
# Helpers de slide
# ══════════════════════════════════════════════════════════════════════════════
def _new_slide(prs, fundo=None):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    for sh in list(sl.shapes):
        sh._element.getparent().remove(sh._element)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = fundo or CINZA_BG
    return sl


def _rect(sl, x, y, w, h, rgb, border_rgb=None):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    if border_rgb:
        sh.line.color.rgb = border_rgb
    else:
        sh.line.fill.background()
    return sh


def _txt(sl, text, x, y, w, h, size=12, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def _header(sl, titulo, subtitulo, cfg):
    _txt(sl, titulo,
         Inches(0.6), Inches(0.28), Inches(14), Inches(0.85),
         size=34, bold=True, color=cfg["cor_principal"])
    if subtitulo:
        _txt(sl, subtitulo,
             Inches(0.6), Inches(1.08), Inches(14), Inches(0.45),
             size=12, color=CINZA_TX)
    _rect(sl, Inches(0.6), Inches(1.02), Inches(18.8), Pt(3), cfg["cor_principal"])
    _txt(sl, cfg["nome"].upper(),
         Inches(14.5), Inches(0.28), Inches(5.0), Inches(0.45),
         size=12, bold=True, color=cfg["cor_principal"], align=PP_ALIGN.RIGHT)
    _txt(sl, cfg["bi_label"],
         Inches(14.5), Inches(0.72), Inches(5.0), Inches(0.35),
         size=10, color=CINZA_TX, align=PP_ALIGN.RIGHT)


def _fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _setup_ax(ax, fc="#F4F6F8"):
    ax.set_facecolor(fc)
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CCCCCC")
    ax.tick_params(colors="#556575", labelsize=9)
    ax.grid(axis="y", color="#DDDDDD", linewidth=0.6, zorder=0)


def _fmt_R(v):
    if pd.isna(v):
        return "—"
    if v >= 1_000_000:
        return f"R$ {v/1_000_000:.2f} Mi"
    if v >= 1_000:
        return f"R$ {v/1_000:.1f} K"
    return f"R$ {v:,.2f}"


def _pct(v):
    if pd.isna(v):
        return "—"
    return f"{v:.1f}%"


def _filtrar(df, col, d0, d1):
    if df.empty or col not in df.columns:
        return df
    d = pd.to_datetime(df[col], errors="coerce")
    return df[(d >= pd.Timestamp(d0)) & (d <= pd.Timestamp(d1))]


def _periodo_label(d0, d1):
    d0t = pd.Timestamp(d0)
    d1t = pd.Timestamp(d1)
    if d0t.year == d1t.year:
        return f"{MESES_PT[d0t.month]}/{d0t.year} — {MESES_PT[d1t.month]}/{d1t.year}"
    return f"{MESES_PT[d0t.month]}/{d0t.year} — {MESES_PT[d1t.month]}/{d1t.year}"


def _mes_label(ts):
    return f"{MESES_PT[ts.month]}/{ts.year}"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ══════════════════════════════════════════════════════════════════════════════
def _slide_capa(prs, d0, d1, cfg):
    sl = _new_slide(prs, fundo=cfg["cor_escura"])

    # Faixa lateral esquerda
    _rect(sl, 0, 0, Inches(7.5), H, cfg["cor_principal"])

    # Icone / emoji grande
    _txt(sl, "💧",
         Inches(0.5), Inches(2.0), Inches(6.5), Inches(2.5),
         size=96, align=PP_ALIGN.CENTER)

    # Nome do cliente
    _txt(sl, cfg["nome"].upper(),
         Inches(0.4), Inches(4.6), Inches(6.7), Inches(1.2),
         size=32, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    # Label BI
    _txt(sl, "Business Intelligence",
         Inches(0.4), Inches(5.65), Inches(6.7), Inches(0.6),
         size=16, color=RGBColor(0xC8, 0xDC, 0xF0), align=PP_ALIGN.CENTER)

    # Bloco direito — titulo
    _txt(sl, "Relatório de Fechamento",
         Inches(8.2), Inches(2.8), Inches(10.8), Inches(1.2),
         size=40, bold=True, color=BRANCO)

    _txt(sl, _periodo_label(d0, d1),
         Inches(8.2), Inches(4.0), Inches(10.8), Inches(0.8),
         size=22, color=RGBColor(0xC8, 0xDC, 0xF0))

    hoje_str = pd.Timestamp.today().strftime("%d/%m/%Y")
    _txt(sl, f"Gerado em {hoje_str}",
         Inches(8.2), Inches(4.8), Inches(10.8), Inches(0.5),
         size=13, color=RGBColor(0x8B, 0xA8, 0xC0), italic=True)

    # Linha decorativa
    _rect(sl, Inches(8.2), Inches(5.5), Inches(10.0), Pt(2), cfg["cor_clara"])

    _txt(sl, "Dados extraídos do BigQuery — atualização diária automática",
         Inches(8.2), Inches(5.7), Inches(10.8), Inches(0.5),
         size=10, color=RGBColor(0x70, 0x90, 0xA8), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Resumo Executivo (6 KPIs)
# ══════════════════════════════════════════════════════════════════════════════
def _slide_resumo(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Resumo Executivo", _periodo_label(d0, d1), cfg)

    fat  = D.get("fat",  pd.DataFrame())
    arr  = D.get("arr",  pd.DataFrame())
    inad = D.get("inad", pd.DataFrame())
    cor  = D.get("cor",  pd.DataFrame())
    rel  = D.get("rel",  pd.DataFrame())

    # ── Calculos ──
    vl_fat = 0.0
    if not fat.empty and "vl_total_faturado" in fat.columns and "dt_ref" in fat.columns:
        sub = _filtrar(fat, "dt_ref", d0, d1)
        vl_fat = sub["vl_total_faturado"].sum()

    vl_arr = 0.0
    if not arr.empty and "vl_arrecadado" in arr.columns and "dt_ref" in arr.columns:
        sub = _filtrar(arr, "dt_ref", d0, d1)
        vl_arr = sub["vl_arrecadado"].sum()

    pct_efic = (vl_arr / vl_fat * 100) if vl_fat > 0 else 0.0

    vl_inad = 0.0
    nr_inad = 0
    if not inad.empty:
        col_v = next((c for c in ["vl_debito","vl_total","vl_saldo"] if c in inad.columns), None)
        if col_v:
            vl_inad = inad[col_v].sum()
        nr_inad = len(inad)

    nr_econ = 0
    if not fat.empty and "nr_economia_agua" in fat.columns and "dt_ref" in fat.columns:
        sub = _filtrar(fat, "dt_ref", d0, d1)
        if not sub.empty:
            nr_econ = int(sub["nr_economia_agua"].iloc[-1]) if "dt_ref" in sub.columns else int(sub["nr_economia_agua"].max())

    nr_cor = 0
    if not cor.empty and "dt_fim_execucao" in cor.columns:
        sub = _filtrar(cor, "dt_fim_execucao", d0, d1)
        nr_cor = len(sub)

    kpis = [
        ("Faturamento",   _fmt_R(vl_fat),      cfg["cor_principal"], "💰"),
        ("Arrecadação",   _fmt_R(vl_arr),       cfg["cor_principal"], "🏦"),
        ("Eficiência",    _pct(pct_efic),        VERDE_C if pct_efic >= 90 else LARANJA, "📈"),
        ("Inadimplência", _fmt_R(vl_inad),      VERMELHO, "⚠️"),
        ("Economias",     f"{nr_econ:,}".replace(",","."), cfg["cor_principal"], "🏠"),
        ("Cortes",        f"{nr_cor:,}".replace(",","."), LARANJA if nr_cor > 0 else VERDE_C, "✂️"),
    ]

    # Grade 3x2
    cols = 3
    card_w = Inches(5.8)
    card_h = Inches(3.4)
    gap_x  = Inches(0.5)
    gap_y  = Inches(0.45)
    start_x = Inches(0.6)
    start_y = Inches(1.55)

    for i, (label, valor, cor_kpi, emoji) in enumerate(kpis):
        col = i % cols
        row = i // cols
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)

        _rect(sl, cx, cy, card_w, card_h, BRANCO, border_rgb=RGBColor(0xDD,0xDD,0xDD))
        # faixa superior colorida
        _rect(sl, cx, cy, card_w, Inches(0.22), cor_kpi)
        _txt(sl, emoji + "  " + label,
             cx + Inches(0.22), cy + Inches(0.32), card_w - Inches(0.3), Inches(0.55),
             size=14, color=CINZA_TX)
        _txt(sl, valor,
             cx + Inches(0.18), cy + Inches(0.95), card_w - Inches(0.3), Inches(1.4),
             size=28, bold=True, color=cor_kpi)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Faturamento
# ══════════════════════════════════════════════════════════════════════════════
def _slide_faturamento(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Faturamento", _periodo_label(d0, d1), cfg)

    fat = D.get("fat", pd.DataFrame())
    if fat.empty or "dt_ref" not in fat.columns or "vl_total_faturado" not in fat.columns:
        _txt(sl, "Dados insuficientes.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    sub = _filtrar(fat, "dt_ref", d0, d1).copy()
    sub["_mes"] = pd.to_datetime(sub["dt_ref"]).dt.to_period("M")
    fat_m = sub.groupby("_mes")["vl_total_faturado"].sum().reset_index()
    fat_m["_mes"] = fat_m["_mes"].dt.to_timestamp()
    fat_m = fat_m.sort_values("_mes")

    # Grafico barras mensais
    fig, axes = plt.subplots(1, 2, figsize=(18, 6.5), facecolor="#F4F6F8")
    fig.subplots_adjust(wspace=0.35)

    ax1 = axes[0]
    _setup_ax(ax1)
    labels = [_mes_label(r) for r in fat_m["_mes"]]
    vals   = fat_m["vl_total_faturado"].values
    bars = ax1.bar(labels, vals, color=cfg["cor_hex"], width=0.6, zorder=3)
    ax1.set_title("Faturamento Mensal", fontsize=12, color=cfg["cor_esc_hex"], pad=8)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.0f}K"))
    for bar, val in zip(bars, vals):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(vals)*0.01,
                 f"R${val/1e3:.0f}K", ha="center", va="bottom", fontsize=8,
                 color=cfg["cor_esc_hex"])
    plt.setp(ax1.get_xticklabels(), rotation=45, ha="right")

    # Pizza componentes
    ax2 = axes[1]
    ax2.set_facecolor("#F4F6F8")
    comp_cols = {
        "Água": "vl_agua",
        "Tarifa Básica": "vl_tarifa_basica",
        "Serviços": "vl_servico",
        "Lixo": "vl_lixo",
    }
    comp_cores = ["#1A6FAD", "#8B5CF6", "#F39C12", "#27AE60"]
    comp_vals = []
    comp_lbls = []
    for lbl, col in comp_cols.items():
        if col in sub.columns:
            v = sub[col].sum()
            if v > 0:
                comp_vals.append(v)
                comp_lbls.append(lbl)

    if comp_vals:
        wedges, texts, autotexts = ax2.pie(
            comp_vals, labels=comp_lbls,
            colors=comp_cores[:len(comp_vals)],
            autopct="%1.1f%%", startangle=90,
            textprops={"fontsize": 9})
        for at in autotexts:
            at.set_color("white")
            at.set_fontsize(8)
        ax2.set_title("Componentes do Faturamento", fontsize=12,
                      color=cfg["cor_esc_hex"], pad=8)
    else:
        ax2.text(0.5, 0.5, "Sem dados de componentes", ha="center", va="center",
                 transform=ax2.transAxes, color=CINZA_TX, fontsize=11)
        ax2.axis("off")

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(0.6), Inches(1.5), Inches(18.8), Inches(7.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Arrecadacao
# ══════════════════════════════════════════════════════════════════════════════
def _slide_arrecadacao(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Arrecadação", _periodo_label(d0, d1), cfg)

    fat = D.get("fat", pd.DataFrame())
    arr = D.get("arr", pd.DataFrame())

    if arr.empty or "dt_ref" not in arr.columns or "vl_arrecadado" not in arr.columns:
        _txt(sl, "Dados insuficientes.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    arr_s = _filtrar(arr, "dt_ref", d0, d1).copy()
    arr_s["_mes"] = pd.to_datetime(arr_s["dt_ref"]).dt.to_period("M")
    arr_m = arr_s.groupby("_mes")["vl_arrecadado"].sum().reset_index()
    arr_m["_mes"] = arr_m["_mes"].dt.to_timestamp()
    arr_m = arr_m.sort_values("_mes")

    fat_m = pd.Series(dtype=float)
    if not fat.empty and "dt_ref" in fat.columns and "vl_total_faturado" in fat.columns:
        fat_s = _filtrar(fat, "dt_ref", d0, d1).copy()
        fat_s["_mes"] = pd.to_datetime(fat_s["dt_ref"]).dt.to_period("M")
        fat_agg = fat_s.groupby("_mes")["vl_total_faturado"].sum()
        fat_m = fat_agg

    fig, ax = plt.subplots(figsize=(18, 7), facecolor="#F4F6F8")
    _setup_ax(ax)

    labels = [_mes_label(r) for r in arr_m["_mes"]]
    x = np.arange(len(labels))
    vals_arr = arr_m["vl_arrecadado"].values

    bars = ax.bar(x, vals_arr, color=cfg["cor_hex"], width=0.55, label="Arrecadado", zorder=3)

    # Linha faturamento se disponivel
    if not fat_m.empty:
        vals_fat = []
        for p in arr_m["_mes"].dt.to_period("M"):
            vals_fat.append(fat_m.get(p, np.nan))
        ax2 = ax.twinx()
        ax2.plot(x, vals_fat, "o--", color="#E74C3C", linewidth=1.8,
                 markersize=5, label="Faturado")
        ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.0f}K"))
        ax2.tick_params(colors="#E74C3C", labelsize=8)
        ax2.spines[["top"]].set_visible(False)
        ax2.legend(loc="upper right", fontsize=8)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=45, ha="right")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.0f}K"))
    ax.set_title("Arrecadação Mensal vs Faturamento", fontsize=13,
                 color=cfg["cor_esc_hex"], pad=8)
    ax.legend(loc="upper left", fontsize=8)

    for bar, val in zip(bars, vals_arr):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(vals_arr)*0.005,
                f"R${val/1e3:.0f}K", ha="center", va="bottom", fontsize=7.5,
                color=cfg["cor_esc_hex"])

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(0.6), Inches(1.5), Inches(18.8), Inches(7.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Inadimplencia
# ══════════════════════════════════════════════════════════════════════════════
def _slide_inadimplencia(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Inadimplência", _periodo_label(d0, d1), cfg)

    inad = D.get("inad", pd.DataFrame())
    if inad.empty:
        _txt(sl, "Sem dados de inadimplência.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    col_v = next((c for c in ["vl_debito","vl_total","vl_saldo"] if c in inad.columns), None)
    if not col_v:
        _txt(sl, "Coluna de valor não encontrada.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    total_inad = inad[col_v].sum()
    nr_inad    = len(inad)

    # KPI card superior
    _rect(sl, Inches(0.6), Inches(1.55), Inches(6), Inches(2.4), cfg["cor_escura"])
    _txt(sl, "Total Inadimplente", Inches(0.9), Inches(1.8), Inches(5.5), Inches(0.6),
         size=13, color=RGBColor(0xC8,0xDC,0xF0))
    _txt(sl, _fmt_R(total_inad), Inches(0.9), Inches(2.3), Inches(5.5), Inches(1.0),
         size=28, bold=True, color=BRANCO)
    _txt(sl, f"{nr_inad:,} registros".replace(",","."), Inches(0.9), Inches(3.25), Inches(5.5), Inches(0.5),
         size=12, color=RGBColor(0x8B,0xA8,0xC0))

    # Aging por faixa (se disponivel)
    col_vc = next((c for c in ["dt_vencimento","dt_ref_documento"] if c in inad.columns), None)
    fig, ax = plt.subplots(figsize=(12, 6.5), facecolor="#F4F6F8")
    _setup_ax(ax)
    ax.grid(axis="x", color="#DDDDDD", linewidth=0.6, zorder=0)
    ax.grid(axis="y", visible=False)

    if col_vc:
        hoje = pd.Timestamp.today()
        inad2 = inad.copy()
        inad2[col_vc] = pd.to_datetime(inad2[col_vc], errors="coerce")
        inad2["dias"] = (hoje - inad2[col_vc]).dt.days
        faixas = [
            ("0–30 dias",    inad2[inad2["dias"].between(0, 30)][col_v].sum()),
            ("31–60 dias",   inad2[inad2["dias"].between(31, 60)][col_v].sum()),
            ("61–90 dias",   inad2[inad2["dias"].between(61, 90)][col_v].sum()),
            ("91–180 dias",  inad2[inad2["dias"].between(91, 180)][col_v].sum()),
            (">180 dias",    inad2[inad2["dias"] > 180][col_v].sum()),
        ]
        lbls_f = [f[0] for f in faixas]
        vals_f = [f[1] for f in faixas]
        cores_f = ["#2ECC71","#F1C40F","#E67E22","#E74C3C","#8E44AD"]
        bars = ax.barh(lbls_f, vals_f, color=cores_f, height=0.55, zorder=3)
        ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.0f}K"))
        for bar, val in zip(bars, vals_f):
            ax.text(bar.get_width() + max(vals_f)*0.01, bar.get_y() + bar.get_height()/2,
                    f"R${val/1e3:.1f}K", va="center", fontsize=9, color="#333333")
        ax.set_title("Aging — Inadimplência por Faixa de Vencimento", fontsize=12,
                     color=cfg["cor_esc_hex"], pad=8)
    else:
        ax.text(0.5, 0.5, "Dados de aging não disponíveis", ha="center", va="center",
                transform=ax.transAxes, fontsize=12, color="#888888")
        ax.axis("off")

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(7.2), Inches(1.5), Inches(12.2), Inches(7.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Economias e Ligacoes
# ══════════════════════════════════════════════════════════════════════════════
def _slide_economias(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Economias e Ligações", _periodo_label(d0, d1), cfg)

    fat = D.get("fat", pd.DataFrame())
    if fat.empty or "dt_ref" not in fat.columns:
        _txt(sl, "Dados insuficientes.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    sub = _filtrar(fat, "dt_ref", d0, d1).copy()
    sub["_mes"] = pd.to_datetime(sub["dt_ref"]).dt.to_period("M")

    cols_eco = {c: c for c in ["nr_economia_agua", "nr_economia_esgoto",
                                "nr_lig_agua", "nr_lig_esgoto"] if c in sub.columns}
    if not cols_eco:
        _txt(sl, "Colunas de economias não disponíveis.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    agg = sub.groupby("_mes")[list(cols_eco.values())].last().reset_index()
    agg["_mes"] = agg["_mes"].dt.to_timestamp()
    agg = agg.sort_values("_mes")
    labels = [_mes_label(r) for r in agg["_mes"]]
    x = np.arange(len(labels))

    fig, ax = plt.subplots(figsize=(18, 7), facecolor="#F4F6F8")
    _setup_ax(ax)

    width = 0.35
    cores_bar = [cfg["cor_hex"], cfg["cor_clr_hex"], "#27AE60", "#8B5CF6"]
    nomes = {"nr_economia_agua": "Econ. Água", "nr_economia_esgoto": "Econ. Esgoto",
             "nr_lig_agua": "Lig. Água", "nr_lig_esgoto": "Lig. Esgoto"}

    n_cols = len(cols_eco)
    offsets = np.linspace(-(n_cols-1)*width/2, (n_cols-1)*width/2, n_cols)
    for i, (col, nome) in enumerate([(c, nomes.get(c, c)) for c in cols_eco]):
        if col in agg.columns:
            ax.bar(x + offsets[i], agg[col], width=width,
                   label=nome, color=cores_bar[i % len(cores_bar)], zorder=3)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=45, ha="right")
    ax.set_title("Economias e Ligações por Mês", fontsize=13, color=cfg["cor_esc_hex"], pad=8)
    ax.legend(fontsize=9)

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(0.6), Inches(1.5), Inches(18.8), Inches(7.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Cortes e Religacoes
# ══════════════════════════════════════════════════════════════════════════════
def _slide_cortes(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Cortes e Religações", _periodo_label(d0, d1), cfg)

    cor = D.get("cor", pd.DataFrame())
    rel = D.get("rel", pd.DataFrame())

    col_cor = next((c for c in ["dt_fim_execucao","dt_execucao","dt_ref"] if cor is not None and c in cor.columns), None)
    col_rel = next((c for c in ["dt_reliagacao","dt_fim_execucao","dt_ref"] if rel is not None and c in rel.columns), None)

    def _mes_agg(df, col):
        if df is None or df.empty or not col:
            return pd.Series(dtype=int)
        s = df.copy()
        s["_mes"] = pd.to_datetime(s[col], errors="coerce").dt.to_period("M")
        s = s[(s[col] >= str(d0)) & (s[col] <= str(d1))]
        return s.groupby("_mes").size()

    cor_m = _mes_agg(cor, col_cor)
    rel_m = _mes_agg(rel, col_rel)

    all_meses = sorted(set(list(cor_m.index) + list(rel_m.index)))
    if not all_meses:
        _txt(sl, "Sem dados de cortes/religações no período.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    labels = [f"{MESES_PT[m.month]}/{m.year}" for m in [p.to_timestamp() for p in all_meses]]
    vals_c = [cor_m.get(m, 0) for m in all_meses]
    vals_r = [rel_m.get(m, 0) for m in all_meses]
    x = np.arange(len(labels))

    fig, ax = plt.subplots(figsize=(18, 7), facecolor="#F4F6F8")
    _setup_ax(ax)

    ax.bar(x - 0.2, vals_c, 0.38, label="Cortes", color="#E74C3C", zorder=3)
    ax.bar(x + 0.2, vals_r, 0.38, label="Religações", color="#27AE60", zorder=3)

    if vals_c:
        ax.axhline(np.mean(vals_c), color="#E74C3C", linewidth=1.2, linestyle="--", alpha=0.7,
                   label=f"Média cortes ({np.mean(vals_c):.0f})")
    if vals_r:
        ax.axhline(np.mean(vals_r), color="#27AE60", linewidth=1.2, linestyle="--", alpha=0.7,
                   label=f"Média rel. ({np.mean(vals_r):.0f})")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=45, ha="right")
    ax.set_title("Cortes e Religações por Mês", fontsize=13, color=cfg["cor_esc_hex"], pad=8)
    ax.legend(fontsize=9)

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(0.6), Inches(1.5), Inches(18.8), Inches(7.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Servicos Operacionais
# ══════════════════════════════════════════════════════════════════════════════
def _slide_servicos(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Serviços Operacionais", _periodo_label(d0, d1), cfg)

    srv = D.get("srv", pd.DataFrame())
    if srv.empty:
        _txt(sl, "Sem dados de serviços.", Inches(0.6), Inches(2), Inches(18), Inches(1),
             size=14, color=CINZA_TX)
        return

    col_dt = next((c for c in ["dt_fim_execucao","dt_abertura","dt_ref"] if c in srv.columns), None)
    if col_dt:
        sub = _filtrar(srv, col_dt, d0, d1)
    else:
        sub = srv

    total = len(sub)

    # SLA (se disponivel)
    pct_sla = None
    if "sla_status" in sub.columns:
        sla_ok = (sub["sla_status"].str.lower() == "ok").sum()
        pct_sla = sla_ok / total * 100 if total > 0 else 0

    # KPI cards
    _rect(sl, Inches(0.6), Inches(1.55), Inches(4.0), Inches(1.8), cfg["cor_principal"])
    _txt(sl, "Total Serviços", Inches(0.85), Inches(1.75), Inches(3.5), Inches(0.5),
         size=11, color=RGBColor(0xC8,0xDC,0xF0))
    _txt(sl, f"{total:,}".replace(",","."), Inches(0.85), Inches(2.1), Inches(3.5), Inches(0.9),
         size=30, bold=True, color=BRANCO)

    if pct_sla is not None:
        cor_sla = VERDE_C if pct_sla >= 90 else (LARANJA if pct_sla >= 75 else VERMELHO)
        _rect(sl, Inches(5.0), Inches(1.55), Inches(4.0), Inches(1.8), cfg["cor_escura"])
        _txt(sl, "% SLA", Inches(5.25), Inches(1.75), Inches(3.5), Inches(0.5),
             size=11, color=RGBColor(0xC8,0xDC,0xF0))
        _txt(sl, f"{pct_sla:.1f}%", Inches(5.25), Inches(2.1), Inches(3.5), Inches(0.9),
             size=30, bold=True, color=BRANCO)

    # Top tipos de servico
    col_tipo = next((c for c in ["nm_servico","ds_servico","id_servico_definicao"] if c in sub.columns), None)
    fig, ax = plt.subplots(figsize=(13, 6.5), facecolor="#F4F6F8")
    _setup_ax(ax)
    ax.grid(axis="x", color="#DDDDDD", linewidth=0.6, zorder=0)
    ax.grid(axis="y", visible=False)

    if col_tipo:
        top = sub[col_tipo].value_counts().head(10)
        if not top.empty:
            cores_top = [cfg["cor_hex"]] * len(top)
            bars = ax.barh(top.index[::-1], top.values[::-1], color=cores_top[::-1],
                           height=0.55, zorder=3)
            for bar, val in zip(bars, top.values[::-1]):
                ax.text(bar.get_width() + max(top.values)*0.01,
                        bar.get_y() + bar.get_height()/2,
                        f"{val}", va="center", fontsize=9, color="#333333")
            ax.set_title("Top 10 Tipos de Serviço", fontsize=12,
                         color=cfg["cor_esc_hex"], pad=8)
    else:
        ax.text(0.5, 0.5, "Detalhamento por tipo não disponível",
                ha="center", va="center", transform=ax.transAxes,
                fontsize=11, color="#888888")
        ax.axis("off")

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(9.8), Inches(1.5), Inches(9.6), Inches(7.5))

    # Evolucao mensal
    if col_dt:
        sub2 = sub.copy()
        sub2["_mes"] = pd.to_datetime(sub2[col_dt], errors="coerce").dt.to_period("M")
        srv_m = sub2.groupby("_mes").size().reset_index(name="qtd")
        srv_m["_mes"] = srv_m["_mes"].dt.to_timestamp()
        srv_m = srv_m.sort_values("_mes")

        fig2, ax2 = plt.subplots(figsize=(8.5, 5), facecolor="#F4F6F8")
        _setup_ax(ax2)
        xlabels = [_mes_label(r) for r in srv_m["_mes"]]
        ax2.plot(xlabels, srv_m["qtd"], "o-", color=cfg["cor_hex"],
                 linewidth=2, markersize=5)
        plt.setp(ax2.get_xticklabels(), rotation=45, ha="right")
        ax2.set_title("Serviços por Mês", fontsize=11, color=cfg["cor_esc_hex"], pad=6)
        buf2 = _fig_to_buf(fig2)
        sl.shapes.add_picture(buf2, Inches(0.6), Inches(3.55), Inches(8.8), Inches(5.2))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Perdas (condicional)
# ══════════════════════════════════════════════════════════════════════════════
def _slide_perdas(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Perdas de Água", _periodo_label(d0, d1), cfg)

    prod = D.get("prod_agua", pd.DataFrame())

    col_dt = next((c for c in ["dt_ref","data"] if not prod.empty and c in prod.columns), None)
    col_vp = next((c for c in ["vl_producao","vl_produzido","volume_produzido"] if not prod.empty and c in prod.columns), None)
    col_vc = next((c for c in ["vl_consumido","volume_consumido","vl_consumo"] if not prod.empty and c in prod.columns), None)

    if not col_dt or (not col_vp and not col_vc):
        _txt(sl, "Dados de produção/perdas não disponíveis no período.",
             Inches(0.6), Inches(2.5), Inches(18), Inches(1), size=14, color=CINZA_TX)
        return

    sub = _filtrar(prod, col_dt, d0, d1).copy()
    sub["_mes"] = pd.to_datetime(sub[col_dt], errors="coerce").dt.to_period("M")
    agg_cols = {}
    if col_vp: agg_cols["produzido"] = (col_vp, "sum")
    if col_vc: agg_cols["consumido"] = (col_vc, "sum")
    agg = sub.groupby("_mes").agg(**agg_cols).reset_index()
    agg["_mes"] = agg["_mes"].dt.to_timestamp()
    agg = agg.sort_values("_mes")

    if "produzido" in agg.columns and "consumido" in agg.columns:
        agg["perda"] = agg["produzido"] - agg["consumido"]
        agg["pct_perda"] = agg["perda"] / agg["produzido"] * 100

    labels = [_mes_label(r) for r in agg["_mes"]]
    x = np.arange(len(labels))

    fig, ax = plt.subplots(figsize=(18, 7), facecolor="#F4F6F8")
    _setup_ax(ax)

    if "produzido" in agg.columns:
        ax.bar(x - 0.18, agg["produzido"], 0.34, label="Produzido (m³)", color=cfg["cor_hex"], zorder=3)
    if "consumido" in agg.columns:
        ax.bar(x + 0.18, agg["consumido"], 0.34, label="Consumido (m³)", color=cfg["cor_clr_hex"], zorder=3)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=45, ha="right")
    ax.set_title("Volume Produzido vs Consumido (m³)", fontsize=13, color=cfg["cor_esc_hex"], pad=8)
    ax.legend(fontsize=9)

    if "pct_perda" in agg.columns:
        ax2 = ax.twinx()
        ax2.plot(x, agg["pct_perda"], "s--", color="#E74C3C", linewidth=1.5,
                 markersize=5, label="% Perda")
        ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"{v:.0f}%"))
        ax2.tick_params(colors="#E74C3C", labelsize=8)
        ax2.spines[["top"]].set_visible(False)
        ax2.legend(loc="upper right", fontsize=8)

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(0.6), Inches(1.5), Inches(18.8), Inches(7.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Energia (condicional)
# ══════════════════════════════════════════════════════════════════════════════
def _slide_energia(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Energia Elétrica", _periodo_label(d0, d1), cfg)

    eng = D.get("energia", pd.DataFrame())
    col_dt = next((c for c in ["dt_ref","Data","data"] if not eng.empty and c in eng.columns), None)
    col_vl = next((c for c in ["vl_total","Valor_Total","vl_energia"] if not eng.empty and c in eng.columns), None)
    col_uc = next((c for c in ["nm_unidade","Unidade","nm_uc"] if not eng.empty and c in eng.columns), None)

    if not col_dt or not col_vl:
        _txt(sl, "Dados de energia não disponíveis no período.",
             Inches(0.6), Inches(2.5), Inches(18), Inches(1), size=14, color=CINZA_TX)
        return

    sub = _filtrar(eng, col_dt, d0, d1).copy()
    sub["_mes"] = pd.to_datetime(sub[col_dt], errors="coerce").dt.to_period("M")
    eng_m = sub.groupby("_mes")[col_vl].sum().reset_index()
    eng_m["_mes"] = eng_m["_mes"].dt.to_timestamp()
    eng_m = eng_m.sort_values("_mes")

    fig, axes = plt.subplots(1, 2, figsize=(18, 6.5), facecolor="#F4F6F8")
    fig.subplots_adjust(wspace=0.3)

    ax1 = axes[0]
    _setup_ax(ax1)
    labels = [_mes_label(r) for r in eng_m["_mes"]]
    bars = ax1.bar(labels, eng_m[col_vl], color="#E67E22", width=0.55, zorder=3)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.0f}K"))
    for bar, val in zip(bars, eng_m[col_vl]):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(eng_m[col_vl])*0.01,
                 f"R${val/1e3:.0f}K", ha="center", va="bottom", fontsize=8, color="#7D4800")
    plt.setp(ax1.get_xticklabels(), rotation=45, ha="right")
    ax1.set_title("Custo Energia por Mês", fontsize=12, color=cfg["cor_esc_hex"], pad=8)

    ax2 = axes[1]
    ax2.set_facecolor("#F4F6F8")
    ax2.spines[["top","right"]].set_visible(False)
    if col_uc:
        top_uc = sub.groupby(col_uc)[col_vl].sum().nlargest(8)
        cores_uc = plt.cm.Oranges(np.linspace(0.4, 0.9, len(top_uc)))
        bars2 = ax2.barh(top_uc.index[::-1], top_uc.values[::-1],
                         color=cores_uc[::-1], height=0.55, zorder=3)
        ax2.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.0f}K"))
        for bar, val in zip(bars2, top_uc.values[::-1]):
            ax2.text(bar.get_width() + max(top_uc.values)*0.01,
                     bar.get_y() + bar.get_height()/2,
                     f"R${val/1e3:.0f}K", va="center", fontsize=8.5, color="#333")
        ax2.grid(axis="x", color="#DDDDDD", linewidth=0.6, zorder=0)
        ax2.tick_params(colors="#556575", labelsize=9)
        ax2.set_title("Top UCs por Custo", fontsize=12, color=cfg["cor_esc_hex"], pad=8)
    else:
        ax2.text(0.5, 0.5, "Dados por UC não disponíveis", ha="center", va="center",
                 transform=ax2.transAxes, fontsize=11, color="#888888")
        ax2.axis("off")

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(0.6), Inches(1.5), Inches(18.8), Inches(7.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Frota (condicional)
# ══════════════════════════════════════════════════════════════════════════════
def _slide_frota(prs, D, d0, d1, cfg):
    sl = _new_slide(prs)
    _header(sl, "Frota — Combustível", _periodo_label(d0, d1), cfg)

    frota = D.get("frota", pd.DataFrame())
    col_dt = next((c for c in ["Data","dt_ref","data"] if not frota.empty and c in frota.columns), None)
    col_vl = next((c for c in ["Valor_Total","vl_total","valor_total"] if not frota.empty and c in frota.columns), None)
    col_lt = next((c for c in ["Quantidade","litros","quantidade"] if not frota.empty and c in frota.columns), None)
    col_mo = next((c for c in ["Motorista","motorista","nm_motorista"] if not frota.empty and c in frota.columns), None)

    if not col_dt or not col_vl:
        _txt(sl, "Dados de frota não disponíveis no período.",
             Inches(0.6), Inches(2.5), Inches(18), Inches(1), size=14, color=CINZA_TX)
        return

    sub = _filtrar(frota, col_dt, d0, d1)
    total_vl = sub[col_vl].sum()
    total_lt  = sub[col_lt].sum() if col_lt else None

    # KPI
    _rect(sl, Inches(0.6), Inches(1.55), Inches(4.5), Inches(1.8), cfg["cor_principal"])
    _txt(sl, "Custo Total", Inches(0.85), Inches(1.75), Inches(4.0), Inches(0.5),
         size=11, color=RGBColor(0xC8,0xDC,0xF0))
    _txt(sl, _fmt_R(total_vl), Inches(0.85), Inches(2.1), Inches(4.0), Inches(0.9),
         size=26, bold=True, color=BRANCO)

    if total_lt:
        _rect(sl, Inches(5.5), Inches(1.55), Inches(3.5), Inches(1.8), cfg["cor_escura"])
        _txt(sl, "Litros", Inches(5.75), Inches(1.75), Inches(3.0), Inches(0.5),
             size=11, color=RGBColor(0xC8,0xDC,0xF0))
        _txt(sl, f"{total_lt:,.0f} L".replace(",","."), Inches(5.75), Inches(2.1), Inches(3.0), Inches(0.9),
             size=26, bold=True, color=BRANCO)

    sub2 = sub.copy()
    sub2["_mes"] = pd.to_datetime(sub2[col_dt], errors="coerce").dt.to_period("M")

    fig, axes = plt.subplots(1, 2, figsize=(14, 5.8), facecolor="#F4F6F8")
    fig.subplots_adjust(wspace=0.35)

    ax1 = axes[0]
    _setup_ax(ax1)
    frota_m = sub2.groupby("_mes")[col_vl].sum().reset_index()
    frota_m["_mes"] = frota_m["_mes"].dt.to_timestamp()
    frota_m = frota_m.sort_values("_mes")
    labels_m = [_mes_label(r) for r in frota_m["_mes"]]
    ax1.bar(labels_m, frota_m[col_vl], color="#3D5A80", width=0.55, zorder=3)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.1f}K"))
    plt.setp(ax1.get_xticklabels(), rotation=45, ha="right")
    ax1.set_title("Custo por Mês", fontsize=11, color=cfg["cor_esc_hex"], pad=6)

    ax2 = axes[1]
    ax2.set_facecolor("#F4F6F8")
    ax2.spines[["top","right"]].set_visible(False)
    if col_mo:
        top_mo = sub.groupby(col_mo)[col_vl].sum().nlargest(8)
        cores_mo = plt.cm.Blues(np.linspace(0.4, 0.9, len(top_mo)))
        ax2.barh(top_mo.index[::-1], top_mo.values[::-1], color=cores_mo[::-1], height=0.55, zorder=3)
        ax2.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"R${v/1e3:.1f}K"))
        ax2.grid(axis="x", color="#DDDDDD", linewidth=0.6, zorder=0)
        ax2.tick_params(colors="#556575", labelsize=9)
        ax2.set_title("Top Motoristas por Custo", fontsize=11, color=cfg["cor_esc_hex"], pad=6)
    else:
        ax2.text(0.5, 0.5, "Dados por motorista não disponíveis",
                 ha="center", va="center", transform=ax2.transAxes, fontsize=10, color="#888888")
        ax2.axis("off")

    buf = _fig_to_buf(fig)
    sl.shapes.add_picture(buf, Inches(0.6), Inches(3.6), Inches(18.8), Inches(5.2))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Encerramento
# ══════════════════════════════════════════════════════════════════════════════
def _slide_encerramento(prs, d0, d1, cfg):
    sl = _new_slide(prs, fundo=cfg["cor_escura"])

    # Faixa lateral
    _rect(sl, 0, 0, Inches(7.5), H, cfg["cor_principal"])

    _txt(sl, "💧",
         Inches(0.5), Inches(2.5), Inches(6.5), Inches(2.5),
         size=96, align=PP_ALIGN.CENTER)

    _txt(sl, cfg["nome"].upper(),
         Inches(0.4), Inches(5.0), Inches(6.7), Inches(1.0),
         size=28, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    _txt(sl, "Obrigado!",
         Inches(8.2), Inches(3.0), Inches(10.8), Inches(1.5),
         size=56, bold=True, color=BRANCO)

    _txt(sl, "Relatório gerado automaticamente pelo",
         Inches(8.2), Inches(4.6), Inches(10.8), Inches(0.6),
         size=14, color=RGBColor(0xC8, 0xDC, 0xF0))

    _txt(sl, f"BI {cfg['nome']}",
         Inches(8.2), Inches(5.1), Inches(10.8), Inches(0.8),
         size=20, bold=True, color=cfg["cor_clara"])

    _txt(sl, _periodo_label(d0, d1),
         Inches(8.2), Inches(5.9), Inches(10.8), Inches(0.6),
         size=14, color=RGBColor(0x8B, 0xA8, 0xC0))

    hoje_str = pd.Timestamp.today().strftime("%d/%m/%Y")
    _txt(sl, f"Emissão: {hoje_str}",
         Inches(8.2), Inches(6.5), Inches(10.8), Inches(0.5),
         size=12, color=RGBColor(0x70, 0x90, 0xA8), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# FUNCAO PRINCIPAL
# ══════════════════════════════════════════════════════════════════════════════
def gerar_pptx(D: dict, d0, d1, cfg: dict) -> bytes:
    """
    Gera a apresentacao PPTX e retorna bytes prontos para download.

    Parametros:
        D   — dicionario de DataFrames (mesmo do app.py load())
        d0  — data inicio (date ou str)
        d1  — data fim   (date ou str)
        cfg — CFG_IPAMERI ou CFG_BURITI

    Retorno:
        bytes do arquivo .pptx
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # 1 — Capa
    _slide_capa(prs, d0, d1, cfg)

    # 2 — Resumo Executivo
    _slide_resumo(prs, D, d0, d1, cfg)

    # 3 — Faturamento
    fat = D.get("fat", pd.DataFrame())
    if not fat.empty:
        _slide_faturamento(prs, D, d0, d1, cfg)

    # 4 — Arrecadacao
    arr = D.get("arr", pd.DataFrame())
    if not arr.empty:
        _slide_arrecadacao(prs, D, d0, d1, cfg)

    # 5 — Inadimplencia
    inad = D.get("inad", pd.DataFrame())
    if not inad.empty:
        _slide_inadimplencia(prs, D, d0, d1, cfg)

    # 6 — Economias e Ligacoes
    if not fat.empty:
        has_eco = any(c in fat.columns for c in ["nr_economia_agua","nr_lig_agua"])
        if has_eco:
            _slide_economias(prs, D, d0, d1, cfg)

    # 7 — Cortes e Religacoes
    cor = D.get("cor", pd.DataFrame())
    rel = D.get("rel", pd.DataFrame())
    if not cor.empty or not rel.empty:
        _slide_cortes(prs, D, d0, d1, cfg)

    # 8 — Servicos
    srv = D.get("srv", pd.DataFrame())
    if not srv.empty:
        _slide_servicos(prs, D, d0, d1, cfg)

    # 9 — Perdas (condicional)
    prod = D.get("prod_agua", pd.DataFrame())
    if prod is not None and not prod.empty:
        _slide_perdas(prs, D, d0, d1, cfg)

    # 10 — Energia (condicional)
    eng = D.get("energia", pd.DataFrame())
    if eng is not None and not eng.empty:
        _slide_energia(prs, D, d0, d1, cfg)

    # 11 — Frota (condicional)
    frota = D.get("frota", pd.DataFrame())
    if frota is not None and not frota.empty:
        _slide_frota(prs, D, d0, d1, cfg)

    # 12 — Encerramento
    _slide_encerramento(prs, d0, d1, cfg)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
