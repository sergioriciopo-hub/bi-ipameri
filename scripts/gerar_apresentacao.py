# -*- coding: utf-8 -*-
"""
Gerador de Apresentacao Executiva - Aguas de Ipameri
Modelo visual baseado em PW Comercial_03.2026.pptx
Slide: 20" x 11.25" | Azul institucional #1A6FAD
"""

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pathlib import Path
from datetime import datetime
import io, warnings
warnings.filterwarnings("ignore")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Configuracoes ─────────────────────────────────────────────────────────────
DATA_DIR   = Path(__file__).parent.parent / "data"
OUTPUT_DIR = Path(__file__).parent.parent / "Apresentacoes"
OUTPUT_DIR.mkdir(exist_ok=True)

# Cores institucionais (adaptadas do template PW Comercial)
AZUL        = RGBColor(0x1A, 0x6F, 0xAD)   # #1A6FAD  — cor principal
AZUL_ESC    = RGBColor(0x0D, 0x3B, 0x5E)   # #0D3B5E  — azul escuro (dados/rodape)
AZUL_CLR    = RGBColor(0x5B, 0x8F, 0xB8)   # #5B8FB8  — azul claro
CINZA_BG    = RGBColor(0xF4, 0xF6, 0xF8)   # fundo dos slides
CINZA_TX    = RGBColor(0x55, 0x65, 0x75)   # texto secundario
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
VERMELHO    = RGBColor(0xC0, 0x39, 0x2B)
VERDE       = RGBColor(0x27, 0xAE, 0x60)
LARANJA     = RGBColor(0xE6, 0x7E, 0x22)

# Slide 20" x 11.25"
W = Inches(20)
H = Inches(11.25)

# Margins para area de conteudo
LEFT   = Inches(1.0)
TOP    = Inches(1.8)
RIGHT  = W - Inches(0.5)
BOTTOM = H - Inches(2.2)

# ── Helpers de leitura ────────────────────────────────────────────────────────
def rd(name):
    import pyarrow.parquet as pq
    import pyarrow as pa
    p = DATA_DIR / f"{name}.parquet"
    if not p.exists():
        return pd.DataFrame()
    tbl = pq.read_table(str(p))
    clean = [pa.field(f.name, f.type) for f in tbl.schema]
    df = tbl.cast(pa.schema(clean)).to_pandas()
    for c in df.select_dtypes(include=["datetime64[ns, UTC]", "datetimetz"]).columns:
        df[c] = pd.to_datetime(df[c], utc=True).dt.tz_localize(None)
    for c in df.select_dtypes(include=["datetime64[ns]"]).columns:
        df[c] = pd.to_datetime(df[c], errors="coerce")
    return df


def fmt_R(v):
    if v >= 1_000_000:
        return f"R$ {v/1_000_000:.2f} Mi"
    elif v >= 1_000:
        return f"R$ {v/1_000:.1f} Mil"
    return f"R$ {v:,.2f}"

def fmt_n(v):
    if v >= 1_000_000:
        return f"{v/1_000_000:.2f}M"
    elif v >= 1_000:
        return f"{v/1_000:.1f}K"
    return f"{int(v):,}"

def mes_nome(ts):
    meses = ["Janeiro","Fevereiro","Marco","Abril","Maio","Junho",
             "Julho","Agosto","Setembro","Outubro","Novembro","Dezembro"]
    return f"{meses[ts.month-1]}/{ts.year}"

def fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


# ── Helpers de slide ──────────────────────────────────────────────────────────
def new_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    sl = prs.slides.add_slide(layout)
    for sh in list(sl.shapes):
        sh._element.getparent().remove(sh._element)
    # Fundo branco/cinza claro
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = CINZA_BG
    return sl


def rect(sl, x, y, w, h, rgb, border=False):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    if border:
        sh.line.color.rgb = rgb
    else:
        sh.line.fill.background()
    return sh


def txt(sl, text, x, y, w, h, size=12, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def titulo_slide(sl, titulo, subtitulo=""):
    """Titulo 36pt azul no topo-esquerdo (estilo PW Comercial)"""
    txt(sl, titulo,
        Inches(0.6), Inches(0.35), Inches(14), Inches(0.85),
        size=36, bold=True, color=AZUL)
    if subtitulo:
        txt(sl, subtitulo,
            Inches(0.6), Inches(1.15), Inches(14), Inches(0.45),
            size=13, color=CINZA_TX)
    # Linha decorativa abaixo do titulo
    rect(sl, Inches(0.6), Inches(1.05), Inches(18.8), Inches(0.04), AZUL)
    # Logo textual canto superior direito
    txt(sl, "AGUAS DE IPAMERI",
        Inches(14.5), Inches(0.35), Inches(5.0), Inches(0.45),
        size=13, bold=True, color=AZUL, align=PP_ALIGN.RIGHT)
    txt(sl, "Business Intelligence",
        Inches(14.5), Inches(0.78), Inches(5.0), Inches(0.35),
        size=10, color=CINZA_TX, align=PP_ALIGN.RIGHT)


def caixa_dados(sl, mes_ref, linhas):
    """
    Caixa de dados no canto inferior esquerdo (estilo template).
    linhas: lista de strings, ex: ["11.245 Ligacoes Ativas", "R$ 1,08 Mi Faturado"]
    """
    box_w = Inches(5.2)
    box_h = Inches(0.55 + len(linhas) * 0.58)
    bx = Inches(0.6)
    by = H - Inches(0.5) - box_h
    rect(sl, bx, by, box_w, box_h, AZUL_ESC)
    # Mes/Ano em destaque
    txt(sl, mes_ref,
        bx + Inches(0.18), by + Inches(0.08),
        box_w - Inches(0.3), Inches(0.45),
        size=15, bold=True, color=BRANCO)
    for i, linha in enumerate(linhas):
        txt(sl, linha,
            bx + Inches(0.18), by + Inches(0.5) + i * Inches(0.52),
            box_w - Inches(0.3), Inches(0.5),
            size=12, color=RGBColor(0xC8,0xDC,0xF0))


def caixa_insight(sl, texto, cor_bg=None):
    """Retangulo de insight na parte inferior (estilo template)"""
    if cor_bg is None:
        cor_bg = AZUL
    bx = Inches(6.0)
    bw = Inches(13.4)
    bh = Inches(1.0)
    by = H - Inches(0.45) - bh
    rect(sl, bx, by, bw, bh, cor_bg)
    txt(sl, texto,
        bx + Inches(0.25), by + Inches(0.12),
        bw - Inches(0.4), bh - Inches(0.2),
        size=11, color=BRANCO, bold=False)


def chart_area():
    """Retorna (x, y, w, h) para a area principal do grafico"""
    return Inches(5.8), Inches(1.65), Inches(13.8), Inches(7.8)


# ── Estilo matplotlib padrao ──────────────────────────────────────────────────
MESES_PT = {1:"Jan",2:"Fev",3:"Mar",4:"Abr",5:"Mai",6:"Jun",
            7:"Jul",8:"Ago",9:"Set",10:"Out",11:"Nov",12:"Dez"}

def setup_ax(ax, facecolor="#F4F6F8"):
    ax.set_facecolor(facecolor)
    ax.spines[["top","right"]].set_visible(False)
    ax.spines[["left","bottom"]].set_color("#CCCCCC")
    ax.tick_params(colors="#556575", labelsize=9)
    ax.grid(axis="y", color="#DDDDDD", linewidth=0.6, zorder=0)


# ══════════════════════════════════════════════════════════════════════════════
# CARREGAR DADOS
# ══════════════════════════════════════════════════════════════════════════════
print("Carregando dados...")
fat   = rd("faturamento")
arr   = rd("arrecadacao")
inad  = rd("pendencia_atual")
cor   = rd("cortes")
rel   = rd("religacoes")
lei   = rd("leituras")
frota = rd("frota_combustivel")

# Normalizar datas
for df_, col_ in [(fat,"dt_ref"),(arr,"dt_ref"),
                  (cor,"dt_fim_execucao"),(rel,"dt_reliagacao")]:
    if not df_.empty and col_ in df_.columns:
        df_[col_] = pd.to_datetime(df_[col_], errors="coerce")
if not inad.empty:
    for c_ in ["dt_ref_documento","dt_vencimento"]:
        if c_ in inad.columns:
            inad[c_] = pd.to_datetime(inad[c_], errors="coerce")
if not lei.empty and "dt_ref" in lei.columns:
    lei["dt_ref"] = pd.to_datetime(lei["dt_ref"], errors="coerce")
if not frota.empty and "Data" in frota.columns:
    frota["Data"] = pd.to_datetime(frota["Data"], errors="coerce")

# Periodo: ultimos 12 meses completos
hoje = pd.Timestamp.today().normalize()
d0_12m = hoje - pd.DateOffset(months=12)

# Faturamento mensal
fat_m = pd.DataFrame()
if not fat.empty and "dt_ref" in fat.columns and "vl_total_faturado" in fat.columns:
    fat["_mes"] = fat["dt_ref"].dt.to_period("M")
    cols_agg = {"vl_fat": ("vl_total_faturado","sum")}
    if "nr_economia_agua" in fat.columns:
        cols_agg["nr_eco"] = ("nr_economia_agua","sum")
    if "nr_lig_agua" in fat.columns:
        cols_agg["nr_lig"] = ("nr_lig_agua","sum")
    if "vl_agua" in fat.columns:
        cols_agg["vl_agua"] = ("vl_agua","sum")
    if "vl_servico" in fat.columns:
        cols_agg["vl_serv"] = ("vl_servico","sum")
    if "vl_servico_basico_agua" in fat.columns:
        cols_agg["vl_tbas"] = ("vl_servico_basico_agua","sum")
    if "vl_lixo" in fat.columns:
        cols_agg["vl_lixo"] = ("vl_lixo","sum")
    fat_m = fat.groupby("_mes").agg(**cols_agg).reset_index()
    fat_m["mes_dt"] = fat_m["_mes"].dt.to_timestamp()
    fat_m = fat_m[fat_m["mes_dt"] >= d0_12m].sort_values("mes_dt").tail(12)

# Arrecadacao mensal
arr_m = pd.DataFrame()
if not arr.empty and "dt_ref" in arr.columns and "vl_total_arrecadado" in arr.columns:
    arr["_mes"] = arr["dt_ref"].dt.to_period("M")
    arr_m = arr.groupby("_mes").agg(vl_arr=("vl_total_arrecadado","sum")).reset_index()
    arr_m["mes_dt"] = arr_m["_mes"].dt.to_timestamp()
    arr_m = arr_m[arr_m["mes_dt"] >= d0_12m].sort_values("mes_dt").tail(12)

# Mes de referencia (ultimo mes disponivel)
mes_ref_ts = fat_m["mes_dt"].iloc[-1] if not fat_m.empty else hoje
mes_ref    = mes_nome(mes_ref_ts)

# KPIs do mes de referencia
vl_fat_mes  = fat_m["vl_fat"].iloc[-1]        if not fat_m.empty else 0
vl_fat_ant  = fat_m["vl_fat"].iloc[-2]        if len(fat_m) >= 2 else vl_fat_mes
vl_arr_mes  = arr_m["vl_arr"].iloc[-1]        if not arr_m.empty else 0
vl_arr_ant  = arr_m["vl_arr"].iloc[-2]        if len(arr_m) >= 2 else vl_arr_mes
nr_eco_mes  = fat_m["nr_eco"].iloc[-1]        if not fat_m.empty and "nr_eco" in fat_m.columns else 0
nr_lig_mes  = fat_m["nr_lig"].iloc[-1]        if not fat_m.empty and "nr_lig" in fat_m.columns else 0
vl_agua_mes = fat_m["vl_agua"].iloc[-1]       if not fat_m.empty and "vl_agua" in fat_m.columns else 0
vl_serv_mes = fat_m["vl_serv"].iloc[-1]       if not fat_m.empty and "vl_serv" in fat_m.columns else 0
vl_tbas_mes = fat_m["vl_tbas"].iloc[-1]       if not fat_m.empty and "vl_tbas" in fat_m.columns else 0
vl_lixo_mes = fat_m["vl_lixo"].iloc[-1]       if not fat_m.empty and "vl_lixo" in fat_m.columns else 0
vl_inad     = inad["vl_divida"].sum()          if not inad.empty and "vl_divida" in inad.columns else 0
efic_mes    = (vl_arr_mes/vl_fat_mes*100)      if vl_fat_mes > 0 else 0
ticket_mes  = (vl_fat_mes/nr_eco_mes)         if nr_eco_mes > 0 else 0
tarifa_mes  = (vl_agua_mes/fat["qt_volume_faturado"].sum()
               if not fat.empty and "qt_volume_faturado" in fat.columns else 0)

# Volume leituras do mes
vol_lid_mes = 0; vol_fat_mes = 0
if not lei.empty and "dt_ref" in lei.columns:
    lei_mes = lei[lei["dt_ref"].dt.to_period("M") == mes_ref_ts.to_period("M")]
    if "qt_volume_lido" in lei_mes.columns:
        vol_lid_mes = int(lei_mes["qt_volume_lido"].sum())
    if "qt_volume_faturado" in lei_mes.columns:
        vol_fat_mes = int(lei_mes["qt_volume_faturado"].sum())

print(f"  Mes referencia : {mes_ref}")
print(f"  Faturamento    : {fmt_R(vl_fat_mes)}")
print(f"  Arrecadacao    : {fmt_R(vl_arr_mes)}")
print(f"  Eficiencia     : {efic_mes:.1f}%")
print(f"  Inadimplencia  : {fmt_R(vl_inad)}")
print(f"  Economias      : {fmt_n(nr_eco_mes)}")


# ══════════════════════════════════════════════════════════════════════════════
# CRIAR APRESENTACAO
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
# Fundo azul escuro
bg = sl.background.fill
bg.solid()
bg.fore_color.rgb = AZUL_ESC

# Faixa azul horizontal no meio
rect(sl, 0, Inches(3.8), W, Inches(0.12), AZUL)
# Faixa clara abaixo
rect(sl, 0, Inches(3.92), W, H - Inches(3.92), RGBColor(0x09,0x28,0x44))

# Textos de capa
txt(sl, "RELATORIO EXECUTIVO",
    Inches(1.2), Inches(1.2), Inches(17.6), Inches(1.3),
    size=52, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
txt(sl, "AGUAS DE IPAMERI",
    Inches(1.2), Inches(2.5), Inches(17.6), Inches(0.9),
    size=30, bold=False, color=AZUL_CLR, align=PP_ALIGN.CENTER)
txt(sl, "Business Intelligence  |  Visao Executiva Comercial",
    Inches(1.2), Inches(4.35), Inches(17.6), Inches(0.7),
    size=18, color=RGBColor(0xA0,0xC4,0xDE), align=PP_ALIGN.CENTER)
txt(sl, f"Referencia: {mes_ref}   |   Gerado em {datetime.now().strftime('%d/%m/%Y')}",
    Inches(1.2), Inches(5.3), Inches(17.6), Inches(0.55),
    size=14, color=RGBColor(0x70,0x9A,0xB8), align=PP_ALIGN.CENTER)
txt(sl, "Confidencial  —  Uso Interno",
    Inches(1.2), Inches(10.4), Inches(17.6), Inches(0.45),
    size=11, color=RGBColor(0x50,0x6A,0x80), align=PP_ALIGN.CENTER)

print("  Slide 1: Capa")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — LIGACOES E ECONOMIAS ATIVAS
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Ligacoes e Economias Ativas",
             f"Evolucao mensal — referencia {mes_ref}")

# Grafico
if not fat_m.empty and "nr_eco" in fat_m.columns:
    mf = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in fat_m["mes_dt"]]
    x  = np.arange(len(mf))

    fig, ax = plt.subplots(figsize=(13, 7.5), facecolor="#F4F6F8")
    setup_ax(ax)
    w = 0.40
    b1 = ax.bar(x - w/2, fat_m["nr_eco"], w, label="Economias",
                color="#1A6FAD", alpha=0.88, zorder=3)
    if "nr_lig" in fat_m.columns:
        b2 = ax.bar(x + w/2, fat_m["nr_lig"], w, label="Ligacoes",
                    color="#5B8FB8", alpha=0.88, zorder=3)
    for bar in b1:
        v = bar.get_height()
        if v > 0:
            ax.text(bar.get_x()+bar.get_width()/2, v+15, f"{int(v):,}",
                    ha="center", va="bottom", fontsize=8, color="#0D3B5E", fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(mf, fontsize=9)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{int(v):,}"))
    ax.legend(fontsize=10, framealpha=0.5)
    ax.set_ylabel("Quantidade", fontsize=10, color="#556575")
    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

caixa_dados(sl, mes_ref, [
    f"{fmt_n(nr_eco_mes)} Economias Ativas",
    f"{fmt_n(nr_lig_mes)} Ligacoes Ativas",
])
caixa_insight(sl,
    f"Em {mes_ref}, foram registradas {fmt_n(nr_eco_mes)} economias e {fmt_n(nr_lig_mes)} "
    f"ligacoes ativas. O crescimento de economias reflete a expansao da rede de abastecimento.")

print("  Slide 2: Ligacoes e Economias")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — VOLUME FATURADO (m3)
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Volume Faturado (m³)",
             f"Consumo mensal de agua — referencia {mes_ref}")

if not fat.empty and "qt_volume_faturado" in fat.columns and "dt_ref" in fat.columns:
    fat["_mes2"] = fat["dt_ref"].dt.to_period("M")
    vol_m = fat.groupby("_mes2")["qt_volume_faturado"].sum().reset_index()
    vol_m["mes_dt"] = vol_m["_mes2"].dt.to_timestamp()
    vol_m = vol_m[vol_m["mes_dt"] >= d0_12m].sort_values("mes_dt").tail(12)
    mf = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in vol_m["mes_dt"]]

    vol_mes_ref = int(vol_m["qt_volume_faturado"].iloc[-1]) if not vol_m.empty else 0
    vol_ant_ref = int(vol_m["qt_volume_faturado"].iloc[-2]) if len(vol_m) >= 2 else vol_mes_ref
    delta_vol   = (vol_mes_ref - vol_ant_ref) / vol_ant_ref * 100 if vol_ant_ref > 0 else 0

    fig, ax = plt.subplots(figsize=(13, 7.5), facecolor="#F4F6F8")
    setup_ax(ax)
    bars = ax.bar(mf, vol_m["qt_volume_faturado"]/1000, color="#1A6FAD", alpha=0.88, zorder=3)
    # linha de media
    media = vol_m["qt_volume_faturado"].mean()/1000
    ax.axhline(media, color="#E67E22", ls="--", lw=1.8, label=f"Media {media:.1f} Mil m3")
    for bar in bars:
        v = bar.get_height()
        if v > 0:
            ax.text(bar.get_x()+bar.get_width()/2, v+0.3, f"{v:.1f}",
                    ha="center", va="bottom", fontsize=8, color="#0D3B5E", fontweight="bold")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{v:.0f} Mil"))
    ax.set_ylabel("Volume em Mil m³", fontsize=10, color="#556575")
    ax.tick_params(axis="x", labelsize=9)
    ax.legend(fontsize=10, framealpha=0.5)
    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"{vol_mes_ref/1000:.1f} Mil m3 Faturados",
        f"{delta_vol:+.1f}% vs mes anterior",
        f"Media 12m: {media:.1f} Mil m3",
    ])
    seta = "acima" if delta_vol >= 0 else "abaixo"
    caixa_insight(sl,
        f"Volume faturado em {mes_ref}: {vol_mes_ref/1000:.1f} mil m3, "
        f"{abs(delta_vol):.1f}% {seta} do mes anterior. "
        f"A media dos ultimos 12 meses e de {media:.1f} mil m3/mes.")

print("  Slide 3: Volume Faturado")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — LEITURAS
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Leituras de Hidrometros",
             f"Distribuicao e qualidade — referencia {mes_ref}")

if not lei.empty and "dt_ref" in lei.columns:
    lei_mes = lei[lei["dt_ref"].dt.to_period("M") == mes_ref_ts.to_period("M")]

    n_total  = len(lei_mes)
    n_critica = int(lei_mes["fl_critica"].sum()) if "fl_critica" in lei_mes.columns else 0
    n_erro    = int(lei_mes["fl_erro_leitura"].sum()) if "fl_erro_leitura" in lei_mes.columns else 0
    n_normal  = n_total - n_critica - n_erro
    pct_crit  = n_critica/n_total*100 if n_total > 0 else 0
    pct_erro  = n_erro/n_total*100 if n_total > 0 else 0

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 7.5), facecolor="#F4F6F8",
                                    gridspec_kw={"width_ratios":[1.2, 1.8]})
    setup_ax(ax1); setup_ax(ax2)

    # Pizza
    labels_ = ["Normais","Criticas","Erros"]
    vals_   = [max(n_normal,0), n_critica, n_erro]
    cores_  = ["#1A6FAD","#E74C3C","#E67E22"]
    ax1.pie(vals_, labels=labels_, colors=cores_, autopct="%1.1f%%",
            startangle=90, textprops={"fontsize":10},
            wedgeprops={"edgecolor":"white","linewidth":1.5})
    ax1.set_title("Distribuicao de Leituras", fontsize=11, color="#0D3B5E", pad=10)

    # Barras mensais por tipo (ultimos 6 meses)
    lei["_mesL"] = lei["dt_ref"].dt.to_period("M")
    ult6 = sorted(lei["_mesL"].unique())[-6:]
    lei6 = lei[lei["_mesL"].isin(ult6)]
    grp = lei6.groupby("_mesL").agg(
        total=("fl_critica","count"),
        critica=("fl_critica","sum") if "fl_critica" in lei.columns else ("fl_critica","count")
    ).reset_index()
    grp["mes_dt"] = grp["_mesL"].dt.to_timestamp()
    mf6 = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in grp["mes_dt"]]
    ax2.bar(mf6, grp["total"], color="#5B8FB8", alpha=0.8, label="Total")
    ax2.bar(mf6, grp["critica"], color="#E74C3C", alpha=0.85, label="Criticas")
    ax2.set_title("Leituras Criticas (6 meses)", fontsize=11, color="#0D3B5E", pad=10)
    ax2.legend(fontsize=9)
    ax2.tick_params(axis="x", labelsize=9)
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{int(v):,}"))

    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"{fmt_n(n_total)} Leituras Realizadas",
        f"{n_critica:,} Criticas ({pct_crit:.1f}%)",
        f"{n_erro:,} com Erro ({pct_erro:.1f}%)",
    ])
    caixa_insight(sl,
        f"Em {mes_ref} foram realizadas {fmt_n(n_total)} leituras. "
        f"{pct_crit:.1f}% foram classificadas como criticas e {pct_erro:.1f}% apresentaram erros, "
        f"indicando oportunidades de melhoria operacional.")

print("  Slide 4: Leituras")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — TICKET MEDIO
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Ticket Medio por Economia",
             f"Faturamento medio mensal — referencia {mes_ref}")

if not fat_m.empty and "nr_eco" in fat_m.columns:
    fat_m2 = fat_m.copy()
    fat_m2["ticket"] = fat_m2["vl_fat"] / fat_m2["nr_eco"]
    mf = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in fat_m2["mes_dt"]]
    ticket_ant = fat_m2["ticket"].iloc[-2] if len(fat_m2) >= 2 else ticket_mes
    delta_tick = (ticket_mes - ticket_ant)/ticket_ant*100 if ticket_ant > 0 else 0
    media_tick = fat_m2["ticket"].mean()

    fig, ax = plt.subplots(figsize=(13, 7.5), facecolor="#F4F6F8")
    setup_ax(ax)
    ax.plot(mf, fat_m2["ticket"], "o-", color="#1A6FAD", lw=2.5, ms=7, zorder=3)
    ax.axhline(media_tick, color="#E67E22", ls="--", lw=1.5,
               label=f"Media {fmt_R(media_tick)}")
    ax.fill_between(range(len(mf)), fat_m2["ticket"], media_tick,
                    where=fat_m2["ticket"] >= media_tick, alpha=0.12, color="#1A6FAD")
    ax.fill_between(range(len(mf)), fat_m2["ticket"], media_tick,
                    where=fat_m2["ticket"] < media_tick, alpha=0.12, color="#E74C3C")
    for i, v in enumerate(fat_m2["ticket"]):
        ax.text(i, v+1.5, fmt_R(v), ha="center", va="bottom", fontsize=8,
                color="#0D3B5E", fontweight="bold")
    ax.set_xticks(range(len(mf)))
    ax.set_xticklabels(mf, fontsize=9)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"R$ {v:.0f}"))
    ax.legend(fontsize=10)
    ax.set_ylabel("R$ / Economia", fontsize=10, color="#556575")
    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"{fmt_R(ticket_mes)} Ticket Medio",
        f"{delta_tick:+.1f}% vs mes anterior",
        f"Media 12m: {fmt_R(media_tick)}",
    ])
    seta2 = "crescimento" if delta_tick >= 0 else "reducao"
    caixa_insight(sl,
        f"Ticket medio em {mes_ref}: {fmt_R(ticket_mes)} por economia, variacao de "
        f"{delta_tick:+.1f}% em relacao ao mes anterior. "
        f"A media dos ultimos 12 meses e de {fmt_R(media_tick)}.")

print("  Slide 5: Ticket Medio")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — TARIFA MEDIA
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Tarifa Media de Agua (R$/m³)",
             f"Evolucao da tarifa unitaria — referencia {mes_ref}")

if not fat.empty and "qt_volume_faturado" in fat.columns and "vl_agua" in fat.columns:
    fat["_mes3"] = fat["dt_ref"].dt.to_period("M")
    tar_m = fat.groupby("_mes3").agg(
        vl_ag=("vl_agua","sum"),
        qt_v=("qt_volume_faturado","sum")
    ).reset_index()
    tar_m["tarifa"] = tar_m["vl_ag"] / tar_m["qt_v"]
    tar_m["mes_dt"] = tar_m["_mes3"].dt.to_timestamp()
    tar_m = tar_m[tar_m["mes_dt"] >= d0_12m].sort_values("mes_dt").tail(12)
    mf = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in tar_m["mes_dt"]]
    tar_mes = tar_m["tarifa"].iloc[-1] if not tar_m.empty else 0
    tar_media = tar_m["tarifa"].mean()

    fig, ax = plt.subplots(figsize=(13, 7.5), facecolor="#F4F6F8")
    setup_ax(ax)
    ax.bar(mf, tar_m["tarifa"], color="#1A6FAD", alpha=0.85, zorder=3)
    ax.axhline(tar_media, color="#E67E22", ls="--", lw=1.5,
               label=f"Media R$ {tar_media:.2f}/m3")
    for i, v in enumerate(tar_m["tarifa"]):
        ax.text(i, v+0.02, f"R$ {v:.2f}", ha="center", va="bottom", fontsize=8,
                color="#0D3B5E", fontweight="bold")
    ax.set_xticks(range(len(mf)))
    ax.set_xticklabels(mf, fontsize=9)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"R$ {v:.2f}"))
    ax.legend(fontsize=10)
    ax.set_ylabel("R$/m³", fontsize=10, color="#556575")
    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"R$ {tar_mes:.2f}/m3 Tarifa Agua",
        f"Media 12m: R$ {tar_media:.2f}/m3",
    ])
    caixa_insight(sl,
        f"A tarifa media de agua em {mes_ref} foi de R$ {tar_mes:.2f}/m3. "
        f"A media dos ultimos 12 meses e de R$ {tar_media:.2f}/m3, "
        f"refletindo os reajustes tarifarios do periodo.")

print("  Slide 6: Tarifa Media")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — FATURAMENTO POR COMPONENTE
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Faturamento por Componente",
             f"Composicao do faturamento — referencia {mes_ref}")

if not fat.empty:
    fat_ref = fat[fat["dt_ref"].dt.to_period("M") == mes_ref_ts.to_period("M")]
    comp = {}
    for col_, nome_ in [("vl_agua","Agua"),("vl_servico_basico_agua","Tarifa Basica"),
                         ("vl_servico","Servicos"),("vl_lixo","Lixo"),
                         ("vl_servico_basico_esgoto","Esgoto")]:
        if col_ in fat_ref.columns:
            v_ = fat_ref[col_].sum()
            if v_ > 0:
                comp[nome_] = v_

    if comp:
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 7.5), facecolor="#F4F6F8",
                                        gridspec_kw={"width_ratios":[1,1.6]})
        setup_ax(ax1); setup_ax(ax2)
        cores_comp = ["#1A6FAD","#5B8FB8","#E67E22","#27AE60","#8B4513","#9B59B6"]
        labels_ = list(comp.keys())
        vals_   = list(comp.values())

        # Pizza
        wedge, texts, autotexts = ax1.pie(
            vals_, labels=labels_, colors=cores_comp[:len(vals_)],
            autopct="%1.1f%%", startangle=90,
            textprops={"fontsize":9},
            wedgeprops={"edgecolor":"white","linewidth":1.5})
        ax1.set_title("Composicao (%)", fontsize=11, color="#0D3B5E", pad=10)

        # Barras horizontais com valores
        y_pos = range(len(labels_))
        bars = ax2.barh(labels_, [v/1000 for v in vals_],
                        color=cores_comp[:len(vals_)], alpha=0.88)
        for bar_ in bars:
            v_ = bar_.get_width()
            ax2.text(v_+max(vals_)/1000*0.01, bar_.get_y()+bar_.get_height()/2,
                     fmt_R(v_*1000), va="center", fontsize=9, color="#0D3B5E", fontweight="bold")
        ax2.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"R$ {v:.0f} Mil"))
        ax2.set_title("Valor por Componente (R$ Mil)", fontsize=11, color="#0D3B5E", pad=10)
        ax2.tick_params(labelsize=9)

        plt.tight_layout(pad=0.4)
        cx, cy, cw, ch = chart_area()
        sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"{fmt_R(vl_fat_mes)} Total Faturado",
        f"{fmt_R(vl_agua_mes)} Agua ({vl_agua_mes/vl_fat_mes*100:.1f}%)" if vl_fat_mes > 0 else "Agua",
        f"{fmt_R(vl_serv_mes)} Servicos",
    ])
    caixa_insight(sl,
        f"Faturamento total em {mes_ref}: {fmt_R(vl_fat_mes)}. "
        f"O componente Agua representa {(vl_agua_mes/vl_fat_mes*100):.1f}% do total. "
        f"Tarifa Basica e Servicos complementam a composicao da receita.")

print("  Slide 7: Faturamento por Componente")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — FATURAMENTO MENSAL (evolucao 12m)
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Faturamento Mensal",
             f"Evolucao dos ultimos 12 meses — referencia {mes_ref}")

if not fat_m.empty:
    mf = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in fat_m["mes_dt"]]
    delta_f = (vl_fat_mes - vl_fat_ant)/vl_fat_ant*100 if vl_fat_ant > 0 else 0

    fig, ax = plt.subplots(figsize=(13, 7.5), facecolor="#F4F6F8")
    setup_ax(ax)
    bars = ax.bar(mf, fat_m["vl_fat"]/1000, color="#1A6FAD", alpha=0.88, zorder=3)
    # destaque ultimo mes
    bars[-1].set_color("#0D3B5E")
    media_f = fat_m["vl_fat"].mean()/1000
    ax.axhline(media_f, color="#E67E22", ls="--", lw=1.5, label=f"Media {fmt_R(media_f*1000)}")
    for bar in bars:
        v = bar.get_height()
        ax.text(bar.get_x()+bar.get_width()/2, v+2, f"{v:.0f}",
                ha="center", va="bottom", fontsize=8, color="#0D3B5E", fontweight="bold")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"R$ {v:.0f} Mil"))
    ax.tick_params(axis="x", labelsize=9)
    ax.legend(fontsize=10)
    ax.set_ylabel("R$ Mil", fontsize=10, color="#556575")
    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"{fmt_R(vl_fat_mes)} Faturado",
        f"{delta_f:+.1f}% vs mes anterior",
        f"Media 12m: {fmt_R(media_f*1000)}",
    ])
    seta3 = "crescimento" if delta_f >= 0 else "queda"
    caixa_insight(sl,
        f"Faturamento de {mes_ref}: {fmt_R(vl_fat_mes)}, variacao de {delta_f:+.1f}% "
        f"em relacao ao mes anterior. Media dos ultimos 12 meses: {fmt_R(media_f*1000)}.")

print("  Slide 8: Faturamento Mensal")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — ARRECADACAO E EFICIENCIA
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Arrecadacao e Eficiencia",
             f"Comparativo faturamento x arrecadacao — referencia {mes_ref}")

if not fat_m.empty and not arr_m.empty:
    df_ea = fat_m.merge(arr_m[["_mes","vl_arr"]], on="_mes", how="inner")
    df_ea["efic"] = df_ea["vl_arr"]/df_ea["vl_fat"]*100
    mf = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in df_ea["mes_dt"]]
    delta_arr2 = (vl_arr_mes - vl_arr_ant)/vl_arr_ant*100 if vl_arr_ant > 0 else 0

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(13, 7.5), facecolor="#F4F6F8",
                                    gridspec_kw={"height_ratios":[2,1]})
    setup_ax(ax1); setup_ax(ax2)

    # Cor das barras de arrecadacao por nivel de eficiencia
    def cor_efic_bar(e):
        if e >= 95:   return "#27AE60"   # verde  >= 95%
        elif e >= 85: return "#E67E22"   # laranja 85-94%
        else:         return "#C0392B"   # vermelho < 85%

    cores_arr = [cor_efic_bar(e) for e in df_ea["efic"]]

    x = np.arange(len(mf)); w = 0.40
    ax1.bar(x - w/2, df_ea["vl_fat"]/1000,  w, label="Faturado",
            color="#1A6FAD", alpha=0.75, zorder=3)
    bars_arr = ax1.bar(x + w/2, df_ea["vl_arr"]/1000, w, label="Arrecadado",
                       color=cores_arr, alpha=0.92, zorder=3)
    ax1.set_xticks(x); ax1.set_xticklabels(mf, fontsize=8)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"R$ {v:.0f} Mil"))
    # Legenda customizada com as faixas de eficiência
    from matplotlib.patches import Patch
    leg_extra = [Patch(color="#27AE60", label="Efic. >= 95%"),
                 Patch(color="#E67E22", label="Efic. 85-94%"),
                 Patch(color="#C0392B", label="Efic. < 85%")]
    leg1 = ax1.legend(fontsize=8, framealpha=0.5)
    ax1.legend(handles=leg_extra, fontsize=8, loc="lower right", framealpha=0.5)
    ax1.add_artist(leg1)
    ax1.set_title("Faturamento vs Arrecadacao — cor indica eficiencia", fontsize=10, color="#0D3B5E")

    ax2.plot(mf, df_ea["efic"], "o-", color="#1A6FAD", lw=2, ms=6, zorder=4)
    ax2.axhline(95,  color="#27AE60", ls="--", lw=1.2, label="Meta 95%")
    ax2.axhline(85,  color="#E67E22", ls=":",  lw=1.0, label="Alerta 85%")
    # Colorir pontos individualmente
    for i, (v, e) in enumerate(zip(df_ea["efic"], df_ea["efic"])):
        ax2.plot(i, v, "o", color=cor_efic_bar(e), ms=9, zorder=5)
    ax2.fill_between(range(len(mf)), df_ea["efic"], 95,
                     where=df_ea["efic"] >= 95, alpha=0.15, color="#27AE60")
    ax2.fill_between(range(len(mf)), df_ea["efic"], 95,
                     where=(df_ea["efic"] >= 85) & (df_ea["efic"] < 95), alpha=0.15, color="#E67E22")
    ax2.fill_between(range(len(mf)), df_ea["efic"], 95,
                     where=df_ea["efic"] < 85,  alpha=0.15, color="#C0392B")
    for i, v in enumerate(df_ea["efic"]):
        ax2.text(i, v+0.6, f"{v:.1f}%", ha="center", va="bottom", fontsize=7.5,
                 color=cor_efic_bar(v), fontweight="bold")
    ax2.set_xticks(range(len(mf))); ax2.set_xticklabels(mf, fontsize=8)
    ax2.legend(fontsize=8, loc="lower right")
    ax2.set_title("Eficiencia de Arrecadacao — verde >= 95% | laranja 85-94% | vermelho < 85%",
                  fontsize=9, color="#0D3B5E")

    plt.tight_layout(pad=0.5)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"{fmt_R(vl_arr_mes)} Arrecadado",
        f"{delta_arr2:+.1f}% vs mes anterior",
        f"Eficiencia: {efic_mes:.1f}%",
    ])
    cor_efic = VERDE if efic_mes >= 95 else LARANJA if efic_mes >= 85 else VERMELHO
    caixa_insight(sl,
        f"Arrecadacao em {mes_ref}: {fmt_R(vl_arr_mes)}, eficiencia de {efic_mes:.1f}%. "
        f"{'Acima da meta de 95%' if efic_mes >= 95 else 'Meta de 95% ainda nao atingida — acoes de cobranca recomendadas'}.",
        cor_bg=cor_efic)

print("  Slide 9: Arrecadacao e Eficiencia")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — INADIMPLENCIA
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Inadimplencia — Carteira de Pendencias",
             "Posicao atual dos documentos em aberto")

if not inad.empty and "vl_divida" in inad.columns:
    inad2 = inad.copy()
    inad2["dias_atraso"] = (pd.Timestamp.today() - inad2["dt_vencimento"]).dt.days.clip(lower=0)
    faixas_ = pd.cut(inad2["dias_atraso"],
                     bins=[0,30,60,90,180,365,99999],
                     labels=["1-30 dias","31-60 dias","61-90 dias",
                             "91-180 dias","181-365 dias","+365 dias"])
    inad2["faixa"] = faixas_
    ag = inad2.groupby("faixa", observed=True)["vl_divida"].sum().reset_index()

    n_docs = len(inad2)
    v_total = inad2["vl_divida"].sum()
    dias_med = inad2["dias_atraso"].mean()

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 7.5), facecolor="#F4F6F8",
                                    gridspec_kw={"width_ratios":[1.4,1]})
    setup_ax(ax1); setup_ax(ax2)

    # Agregar também contagem de clientes por faixa
    ag_cnt = inad2.groupby("faixa", observed=True).agg(
        vl_divida=("vl_divida","sum"),
        n_clientes=("vl_divida","count")
    ).reset_index()

    cores_fx = ["#27AE60","#F1C40F","#E67E22","#E74C3C","#C0392B","#8E1010"]
    # Gráfico principal: barras duplas (valor R$ + qtd clientes)
    x_fx = np.arange(len(ag_cnt))
    ax1_twin = ax1.twiny()
    bars_v = ax1.barh(x_fx - 0.2, ag_cnt["vl_divida"]/1000,
                      0.38, color=cores_fx[:len(ag_cnt)], alpha=0.88, label="Valor (R$ Mil)")
    bars_c = ax1_twin.barh(x_fx + 0.2, ag_cnt["n_clientes"],
                           0.38, color=cores_fx[:len(ag_cnt)], alpha=0.40, hatch="//")
    for bar_ in bars_v:
        v_ = bar_.get_width()
        ax1.text(v_+max(ag_cnt["vl_divida"]/1000)*0.01, bar_.get_y()+bar_.get_height()/2,
                 f"R$ {v_:.0f} Mil", va="center", fontsize=8, color="#0D3B5E", fontweight="bold")
    for bar_ in bars_c:
        v_ = bar_.get_width()
        ax1_twin.text(v_+max(ag_cnt["n_clientes"])*0.01, bar_.get_y()+bar_.get_height()/2,
                      f"{int(v_)} cli.", va="center", fontsize=8, color="#556575")
    ax1.set_yticks(x_fx)
    ax1.set_yticklabels(ag_cnt["faixa"].astype(str), fontsize=9)
    ax1.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"R$ {v:.0f} Mil"))
    ax1_twin.set_xlabel("Clientes", fontsize=8, color="#556575")
    ax1.set_title("Valor e Clientes por Faixa de Atraso", fontsize=11, color="#0D3B5E", pad=10)
    ax1.tick_params(labelsize=8)

    # pizza por valor
    ax2.pie(ag_cnt["vl_divida"], labels=ag_cnt["faixa"].astype(str),
            colors=cores_fx[:len(ag_cnt)], autopct="%1.1f%%",
            startangle=90, textprops={"fontsize":8},
            wedgeprops={"edgecolor":"white","linewidth":1.2})
    ax2.set_title("Composicao por Valor (%)", fontsize=11, color="#0D3B5E", pad=10)

    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    pct_fat2 = v_total/vl_fat_mes*100 if vl_fat_mes > 0 else 0
    caixa_dados(sl, "Posicao Atual", [
        f"{fmt_R(v_total)} em Aberto",
        f"{fmt_n(n_docs)} Documentos",
        f"{dias_med:.0f} dias (media)",
        f"{pct_fat2:.1f}% do fat. mensal",
    ])
    caixa_insight(sl,
        f"Carteira de inadimplencia: {fmt_R(v_total)} em {fmt_n(n_docs)} documentos. "
        f"Representa {pct_fat2:.1f}% do faturamento mensal de referencia. "
        f"Atraso medio de {dias_med:.0f} dias — foco em faixas acima de 90 dias.",
        cor_bg=VERMELHO)

print("  Slide 10: Inadimplencia")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — CORTES E RELIGACOES
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Cortes e Religacoes de Fornecimento",
             f"Evolucao dos ultimos 12 meses — referencia {mes_ref}")

if not cor.empty and not rel.empty:
    cor["_mes"] = cor["dt_fim_execucao"].dt.to_period("M")
    rel["_mes"] = rel["dt_reliagacao"].dt.to_period("M")
    cor_m2 = cor[cor["dt_fim_execucao"] >= d0_12m].groupby("_mes").size().reset_index(name="cortes")
    rel_m2 = rel[rel["dt_reliagacao"]   >= d0_12m].groupby("_mes").size().reset_index(name="religs")
    df_cr2 = cor_m2.merge(rel_m2, on="_mes", how="outer").fillna(0).sort_values("_mes").tail(12)
    df_cr2["mes_dt"] = df_cr2["_mes"].dt.to_timestamp()
    mf = [f"{MESES_PT[d.month]}/{str(d.year)[2:]}" for d in df_cr2["mes_dt"]]
    tot_cor = int(df_cr2["cortes"].sum())
    tot_rel = int(df_cr2["religs"].sum())
    taxa_rel = (tot_rel/tot_cor*100) if tot_cor > 0 else 0

    fig, ax = plt.subplots(figsize=(13, 7.5), facecolor="#F4F6F8")
    setup_ax(ax)
    x = np.arange(len(mf)); w = 0.40
    ax.bar(x - w/2, df_cr2["cortes"], w, label="Cortes",    color="#E74C3C", alpha=0.88, zorder=3)
    ax.bar(x + w/2, df_cr2["religs"], w, label="Religacoes", color="#27AE60", alpha=0.88, zorder=3)
    ax.set_xticks(x); ax.set_xticklabels(mf, fontsize=9)
    ax.legend(fontsize=10)
    ax.set_ylabel("Quantidade", fontsize=10, color="#556575")
    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, f"Acum. 12m ({mes_ref})", [
        f"{fmt_n(tot_cor)} Cortes Executados",
        f"{fmt_n(tot_rel)} Religacoes",
        f"{taxa_rel:.1f}% Taxa de Religacao",
    ])
    caixa_insight(sl,
        f"Nos ultimos 12 meses: {fmt_n(tot_cor)} cortes executados e {fmt_n(tot_rel)} religacoes. "
        f"Taxa de religacao de {taxa_rel:.1f}%, indicando que a maioria dos clientes regulariza "
        f"a situacao apos o corte.")

print("  Slide 11: Cortes e Religacoes")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — FROTA E COMBUSTIVEL
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
titulo_slide(sl, "Frota e Combustivel",
             f"Consumo e eficiencia — referencia {mes_ref}")

if not frota.empty and "Data" in frota.columns:
    fr_ref = frota[frota["Data"].dt.to_period("M") == mes_ref_ts.to_period("M")]
    tot_litros  = fr_ref["Quantidade"].sum()   if "Quantidade"   in fr_ref.columns else 0
    tot_custo   = fr_ref["Valor_Total"].sum()  if "Valor_Total"  in fr_ref.columns else 0
    tot_km      = fr_ref["Km_Rodados"].sum()   if "Km_Rodados"   in fr_ref.columns else 0
    efic_frota  = (tot_km/tot_litros)          if tot_litros > 0 else 0

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 7.5), facecolor="#F4F6F8",
                                    gridspec_kw={"width_ratios":[1.3,1]})
    setup_ax(ax1); setup_ax(ax2)

    # Consumo por veiculo
    if "Veiculo" in fr_ref.columns and "Quantidade" in fr_ref.columns:
        v_grp = fr_ref.groupby("Veiculo")["Quantidade"].sum().sort_values(ascending=False).head(10)
        ax1.barh(v_grp.index, v_grp.values, color="#1A6FAD", alpha=0.88)
        for i, (v_) in enumerate(v_grp.values):
            ax1.text(v_+0.3, i, f"{v_:.0f} L", va="center", fontsize=9, color="#0D3B5E")
        ax1.set_title("Consumo por Veiculo (L)", fontsize=11, color="#0D3B5E", pad=10)
        ax1.tick_params(labelsize=9)

    # Eficiencia por veiculo (km/L)
    if "Veiculo" in fr_ref.columns and "Km_Por_Litro" in fr_ref.columns:
        ef_grp = fr_ref.groupby("Veiculo")["Km_Por_Litro"].mean().sort_values(ascending=False).head(10)
        cores_ef = ["#27AE60" if v >= 7 else "#E67E22" if v >= 5 else "#E74C3C"
                    for v in ef_grp.values]
        ax2.barh(ef_grp.index, ef_grp.values, color=cores_ef, alpha=0.88)
        media_ef = ef_grp.mean()
        ax2.axvline(media_ef, color="#E67E22", ls="--", lw=1.5)
        ax2.set_title("Eficiencia (km/L) por Veiculo", fontsize=11, color="#0D3B5E", pad=10)
        ax2.tick_params(labelsize=9)
        ax2.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{v:.1f} km/L"))

    plt.tight_layout(pad=0.4)
    cx, cy, cw, ch = chart_area()
    sl.shapes.add_picture(fig_to_buf(fig), cx, cy, cw, ch)

    caixa_dados(sl, mes_ref, [
        f"{tot_litros:.0f} Litros Consumidos",
        f"{fmt_R(tot_custo)} Gasto",
        f"{tot_km:,.0f} km Rodados",
        f"{efic_frota:.1f} km/Litro",
    ])
    caixa_insight(sl,
        f"Frota em {mes_ref}: {tot_litros:.0f} litros consumidos, custo total {fmt_R(tot_custo)}, "
        f"{tot_km:,.0f} km percorridos. Eficiencia media de {efic_frota:.1f} km/L. "
        f"Custo medio por km: {fmt_R(tot_custo/tot_km) if tot_km > 0 else 'N/D'}.")

print("  Slide 12: Frota e Combustivel")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — ENCERRAMENTO
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide(prs)
bg_ = sl.background.fill
bg_.solid()
bg_.fore_color.rgb = AZUL_ESC

rect(sl, 0, Inches(4.2), W, Inches(0.1), AZUL)
rect(sl, 0, Inches(4.3), W, H - Inches(4.3), RGBColor(0x09, 0x28, 0x44))

txt(sl, "OBRIGADO",
    Inches(1), Inches(1.4), Inches(18), Inches(1.4),
    size=56, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
txt(sl, "Aguas de Ipameri  —  Business Intelligence",
    Inches(1), Inches(2.9), Inches(18), Inches(0.8),
    size=22, color=AZUL_CLR, align=PP_ALIGN.CENTER)

rect(sl, Inches(6.5), Inches(4.8), Inches(7), Inches(0.06), AZUL)

resumo = (f"Faturado: {fmt_R(vl_fat_mes)}   |   Arrecadado: {fmt_R(vl_arr_mes)}   |   "
          f"Eficiencia: {efic_mes:.1f}%   |   Inadimplencia: {fmt_R(vl_inad)}")
txt(sl, resumo,
    Inches(1), Inches(5.1), Inches(18), Inches(0.6),
    size=14, color=RGBColor(0xA0,0xC4,0xDE), align=PP_ALIGN.CENTER)
txt(sl, f"Referencia: {mes_ref}   |   Gerado em {datetime.now().strftime('%d/%m/%Y as %H:%M')}",
    Inches(1), Inches(5.9), Inches(18), Inches(0.5),
    size=12, color=RGBColor(0x70,0x9A,0xB8), align=PP_ALIGN.CENTER)
txt(sl, "Confidencial  —  Uso Interno",
    Inches(1), Inches(10.5), Inches(18), Inches(0.4),
    size=11, color=RGBColor(0x50,0x6A,0x80), align=PP_ALIGN.CENTER)

print("  Slide 13: Encerramento")


# ── Salvar ────────────────────────────────────────────────────────────────────
nome = f"Relatorio_Executivo_Aguas_Ipameri_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
out  = OUTPUT_DIR / nome
prs.save(str(out))

print(f"\nApresentacao gerada com sucesso!")
print(f"  Arquivo : {out}")
print(f"  Slides  : 13")
print(f"  Tamanho : {out.stat().st_size/1024:.0f} KB")
